import os
import xml.etree.ElementTree as ET
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

def estrai_pptx(nome_file_pptx, nome_file_xml):
    if not os.path.exists(nome_file_pptx):
        print(f"Errore: Il file {nome_file_pptx} non esiste.")
        return

    print(f"Analisi di {nome_file_pptx} in corso... (potrebbe volerci un po' data la dimensione)")
    prs = Presentation(nome_file_pptx)
    
    root = ET.Element("Presentazione")
    root.set("nome", nome_file_pptx)

    for i, slide in enumerate(prs.slides):
        slide_node = ET.SubElement(root, "Slide")
        slide_node.set("numero", str(i + 1))
        
        # Note dell'oratore
        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text
            notes_node = ET.SubElement(slide_node, "Note")
            notes_node.text = notes

        # Oggetti nella slide
        oggetti_node = ET.SubElement(slide_node, "Oggetti")
        
        for shape in slide.shapes:
            obj_node = ET.SubElement(oggetti_node, "Oggetto")
            obj_node.set("id", str(shape.shape_id))
            obj_node.set("tipo", str(shape.shape_type))
            
            # Posizione e dimensione
            pos_node = ET.SubElement(obj_node, "Trasformazione")
            pos_node.set("x", str(shape.left))
            pos_node.set("y", str(shape.top))
            pos_node.set("w", str(shape.width))
            pos_node.set("h", str(shape.height))

            # Testo (se presente)
            if shape.has_text_frame:
                text_node = ET.SubElement(obj_node, "Testo")
                text_node.text = shape.text_frame.text

            # Gestione Immagini/Video (Segnaposto)
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                img_node = ET.SubElement(obj_node, "Media")
                img_node.set("tipo", "Immagine")
                img_node.text = "[Segnaposto Immagine]"
            elif shape.shape_type == MSO_SHAPE_TYPE.MEDIA:
                vid_node = ET.SubElement(obj_node, "Media")
                vid_node.set("tipo", "Video")
                vid_node.text = "[Segnaposto Video]"

    # Scrittura del file XML
    tree = ET.ElementTree(root)
    ET.indent(tree, space="  ", level=0)
    tree.write(nome_file_xml, encoding="utf-8", xml_declaration=True)
    print(f"Estrazione completata! Creato il file: {nome_file_xml}")

if __name__ == "__main__":
    # Usiamo il nome esatto del tuo file
    estrai_pptx("Imperio VII.pptx", "struttura_slide.xml")
